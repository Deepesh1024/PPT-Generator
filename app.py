from flask import Flask, render_template, request, jsonify, send_from_directory, url_for
from flask_cors import CORS
import os
import json
from datetime import datetime
import uuid
from langchain_groq import ChatGroq
from huggingface_hub import InferenceClient
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv
from PIL import Image
import io
import traceback

load_dotenv()

app = Flask(__name__)
CORS(app)

# Configuration
GROQ_API_KEY = os.getenv('GROQ_API_KEY')
HF_TOKEN = os.getenv('HF_TOKEN')
UPLOAD_FOLDER = 'static/generated'
IMAGES_FOLDER = 'static/images'
PPT_FOLDER = 'static/presentations'

# Create directories if they don't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(IMAGES_FOLDER, exist_ok=True)
os.makedirs(PPT_FOLDER, exist_ok=True)

# Initialize LangChain Groq
llm = ChatGroq(
    groq_api_key=GROQ_API_KEY,
    model="moonshotai/kimi-k2-instruct-0905",
    temperature=0.7
)

# Initialize Hugging Face Inference Client
print("Initializing image generation client...")
try:
    client = InferenceClient(api_key=HF_TOKEN)
    print("Image generation client initialized successfully!")
except Exception as e:
    print(f"Error initializing image client: {e}")
    client = None

class PPTBot:
    def __init__(self):
        self.structure_bot = StructureBot()
        self.content_bot = ContentBot()
        
    def generate_presentation(self, topic, slides_count=5):
        try:
            # Step 1: Generate slide structure
            print("Generating slide structure...")
            slide_structure = self.structure_bot.create_slide_structure(topic, slides_count)
            
            # Step 2: Generate content for each slide
            print("Generating content for slides...")
            slides_data = []
            
            for i, slide_info in enumerate(slide_structure):
                content = self.content_bot.generate_slide_content(slide_info, topic)
                image_prompt = self.content_bot.generate_image_prompt(slide_info, topic)
                
                slides_data.append({
                    'slide_number': i + 1,
                    'title': slide_info['title'],
                    'design_type': slide_info['design_type'],
                    'content': content,
                    'image_prompt': image_prompt,
                    'image_path': None
                })
            
            # Step 3: Generate images
            print("Generating images...")
            for slide_data in slides_data:
                print(f"Processing slide {slide_data['slide_number']} of {len(slides_data)}")
                if client:
                    image_path = self.generate_image(slide_data['image_prompt'], slide_data['slide_number'])
                    slide_data['image_path'] = image_path
                    print(f"Image path for slide {slide_data['slide_number']}: {image_path}")
            
            # Step 4: Create PowerPoint
            print("Creating PowerPoint presentation...")
            ppt_path = self.create_powerpoint(slides_data, topic)
            
            return {
                'success': True,
                'slides_data': slides_data,
                'ppt_path': ppt_path,
                'message': 'Presentation generated successfully!'
            }
            
        except Exception as e:
            print(f"Error in generate_presentation: {e}")
            traceback.print_exc()
            return {
                'success': False,
                'error': str(e),
                'message': 'Error generating presentation'
            }
    
    def generate_image(self, prompt, slide_number):
        try:
            if not client:
                print("Image client not available")
                return None
                
            # Generate image with Hugging Face Inference Client
            enhanced_prompt = f"{prompt}, high quality, professional, clean design, presentation style, no text"
            
            print(f"Generating image for slide {slide_number} with prompt: {enhanced_prompt[:100]}...")
            
            # Generate image using HF client
            image = client.text_to_image(
                enhanced_prompt,
                model="black-forest-labs/FLUX.1-dev",
            )
            
            # Verify image is valid
            if not image:
                print(f"No image generated for slide {slide_number}")
                return None
            
            # Save image
            image_filename = f"slide_{slide_number}_{uuid.uuid4().hex[:8]}.png"
            image_path = os.path.join(IMAGES_FOLDER, image_filename)
            
            # Ensure directory exists
            os.makedirs(os.path.dirname(image_path), exist_ok=True)
            
            # Save PIL image
            image.save(image_path, "PNG")
            
            # Verify file was saved
            if os.path.exists(image_path):
                file_size = os.path.getsize(image_path)
                print(f"Image saved successfully: {image_filename} (Size: {file_size} bytes)")
                return f"images/{image_filename}"
            else:
                print(f"Failed to save image for slide {slide_number}")
                return None
            
        except Exception as e:
            print(f"Error generating image for slide {slide_number}: {e}")
            traceback.print_exc()
            return None
    
    def create_powerpoint(self, slides_data, topic):
        try:
            # Create presentation
            prs = Presentation()
            
            # Set slide size to widescreen
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            for i, slide_data in enumerate(slides_data):
                print(f"Creating slide {i+1} of {len(slides_data)}: {slide_data['title']}")
                
                # Add slide with blank layout
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(15), Inches(1.5)
                )
                title_frame = title_shape.text_frame
                title_frame.text = slide_data['title']
                
                # Format title
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.font.size = Pt(36)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(31, 78, 121)
                title_paragraph.alignment = PP_ALIGN.CENTER
                
                # Add content based on design type
                try:
                    if slide_data['design_type'] == 'title':
                        self.add_title_slide_content(slide, slide_data)
                    elif slide_data['design_type'] == 'content':
                        self.add_content_slide(slide, slide_data)
                    elif slide_data['design_type'] == 'image_focus':
                        self.add_image_focus_slide(slide, slide_data)
                    elif slide_data['design_type'] == 'conclusion':
                        self.add_conclusion_slide(slide, slide_data)
                    else:
                        self.add_default_slide(slide, slide_data)
                    
                    print(f"Successfully created slide {i+1}")
                    
                except Exception as slide_error:
                    print(f"Error creating slide {i+1}: {slide_error}")
                    traceback.print_exc()
                    # Continue with next slide instead of failing completely
                    continue
            
            # Save presentation
            ppt_filename = f"{topic.replace(' ', '_').replace('/', '_')}_{uuid.uuid4().hex[:8]}.pptx"
            ppt_path = os.path.join(PPT_FOLDER, ppt_filename)
            
            # Ensure directory exists
            os.makedirs(os.path.dirname(ppt_path), exist_ok=True)
            
            prs.save(ppt_path)
            
            # Verify file was saved
            if os.path.exists(ppt_path):
                file_size = os.path.getsize(ppt_path)
                print(f"PowerPoint saved successfully: {ppt_filename} (Size: {file_size} bytes)")
                return f"presentations/{ppt_filename}"
            else:
                print(f"Failed to save PowerPoint: {ppt_path}")
                return None
            
        except Exception as e:
            print(f"Error creating PowerPoint: {e}")
            traceback.print_exc()
            return None
    
    def add_image_to_slide(self, slide, image_path, left, top, width, height):
        """Helper method to safely add images to slides with better error handling"""
        try:
            if not image_path:
                print("No image path provided")
                return False
                
            # Convert relative path to absolute path
            if not os.path.isabs(image_path):
                full_image_path = os.path.join('static', image_path)
            else:
                full_image_path = image_path
            
            # Check if file exists
            if not os.path.exists(full_image_path):
                print(f"Image file does not exist: {full_image_path}")
                return False
            
            # Check file size
            file_size = os.path.getsize(full_image_path)
            if file_size == 0:
                print(f"Image file is empty: {full_image_path}")
                return False
            
            # Verify it's a valid image
            try:
                with Image.open(full_image_path) as img:
                    img.verify()
            except Exception as img_error:
                print(f"Invalid image file: {full_image_path}, Error: {img_error}")
                return False
            
            # Add picture to slide
            slide.shapes.add_picture(full_image_path, left, top, width, height)
            print(f"Successfully added image: {full_image_path}")
            return True
            
        except Exception as e:
            print(f"Error adding image to slide: {e}")
            traceback.print_exc()
            return False
    
    def add_title_slide_content(self, slide, slide_data):
        # Add subtitle
        subtitle_shape = slide.shapes.add_textbox(
            Inches(2), Inches(2.5), Inches(12), Inches(1)
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = slide_data['content'][:100] + "..." if len(slide_data['content']) > 100 else slide_data['content']
        
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(24)
        subtitle_paragraph.font.color.rgb = RGBColor(89, 89, 89)
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add image if available
        if slide_data['image_path']:
            self.add_image_to_slide(
                slide, slide_data['image_path'],
                Inches(6), Inches(4), Inches(4), Inches(3)
            )
    
    def add_content_slide(self, slide, slide_data):
        # Split content into left text and right image
        content_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(7), Inches(6)
        )
        content_frame = content_shape.text_frame
        content_frame.text = slide_data['content']
        content_frame.word_wrap = True
        
        # Format content
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
        
        # Add image on the right
        if slide_data['image_path']:
            self.add_image_to_slide(
                slide, slide_data['image_path'],
                Inches(8.5), Inches(2), Inches(6.5), Inches(5)
            )
    
    def add_image_focus_slide(self, slide, slide_data):
        # Large image with minimal text
        if slide_data['image_path']:
            self.add_image_to_slide(
                slide, slide_data['image_path'],
                Inches(2), Inches(2), Inches(12), Inches(6)
            )
        
        # Add caption below image
        caption_shape = slide.shapes.add_textbox(
            Inches(1), Inches(8.2), Inches(14), Inches(0.8)
        )
        caption_frame = caption_shape.text_frame
        caption_text = slide_data['content'][:150] + "..." if len(slide_data['content']) > 150 else slide_data['content']
        caption_frame.text = caption_text
        
        caption_paragraph = caption_frame.paragraphs[0]
        caption_paragraph.font.size = Pt(14)
        caption_paragraph.font.italic = True
        caption_paragraph.font.color.rgb = RGBColor(102, 102, 102)
        caption_paragraph.alignment = PP_ALIGN.CENTER
    
    def add_conclusion_slide(self, slide, slide_data):
        # Centered content
        content_shape = slide.shapes.add_textbox(
            Inches(2), Inches(2.5), Inches(12), Inches(4)
        )
        content_frame = content_shape.text_frame
        content_frame.text = slide_data['content']
        
        # Format content
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(31, 78, 121)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Add image if available
        if slide_data['image_path']:
            self.add_image_to_slide(
                slide, slide_data['image_path'],
                Inches(6), Inches(6.5), Inches(4), Inches(2)
            )
        
        # Add decorative shape
        try:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(7), Inches(7), Inches(2), Inches(0.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(31, 78, 121)
        except Exception as e:
            print(f"Error adding decorative shape: {e}")
    
    def add_default_slide(self, slide, slide_data):
        self.add_content_slide(slide, slide_data)

class StructureBot:
    def create_slide_structure(self, topic, slides_count):
        prompt = f"""
        Create a structure for a {slides_count}-slide PowerPoint presentation on "{topic}".
        
        For each slide, provide:
        1. Title (concise and engaging)
        2. Design type (choose from: title, content, image_focus, conclusion)
        3. Key points to cover
        
        Make the presentation flow logically with:
        - Slide 1: Always title slide introducing the topic
        - Middle slides: Mix of content and image_focus slides
        - Last slide: Always conclusion slide
        
        Return ONLY a valid JSON array with this exact format:
        [
            {{
                "slide_number": 1,
                "title": "Slide Title",
                "design_type": "title",
                "key_points": ["point1", "point2", "point3"]
            }}
        ]
        
        Do not include any text before or after the JSON.
        """
        
        try:
            response = llm.invoke(prompt)
            structure_text = response.content.strip()
            
            # Clean up the response to extract JSON
            if "```":
                structure_text = structure_text.split("```json").split("```")
            elif "```" in structure_text:
                # Find content between first and last ```
                parts = structure_text.split("```")
                if len(parts) >= 3:
                    structure_text = parts[1]
                else:
                    structure_text = parts[-1]
            
            # Remove any leading/trailing whitespace
            structure_text = structure_text.strip()
            
            # Parse JSON
            structure = json.loads(structure_text)
            
            # Validate structure
            if not isinstance(structure, list):
                raise ValueError("Structure must be a list")
                
            for slide in structure:
                if not all(key in slide for key in ['slide_number', 'title', 'design_type', 'key_points']):
                    raise ValueError("Each slide must have required keys")
            
            return structure
            
        except Exception as e:
            print(f"Error creating slide structure: {e}")
            print(f"Raw response: {response.content if 'response' in locals() else 'No response'}")
            
            # Fallback structure
            fallback_structure = []
            for i in range(slides_count):
                if i == 0:
                    slide = {
                        "slide_number": 1,
                        "title": f"Introduction to {topic}",
                        "design_type": "title",
                        "key_points": ["Overview", "Objectives", "Key Topics"]
                    }
                elif i == slides_count - 1:
                    slide = {
                        "slide_number": i + 1,
                        "title": "Conclusion",
                        "design_type": "conclusion",
                        "key_points": ["Summary", "Key Takeaways", "Next Steps"]
                    }
                else:
                    design_type = "content" if i % 2 == 1 else "image_focus"
                    slide = {
                        "slide_number": i + 1,
                        "title": f"{topic} - Part {i}",
                        "design_type": design_type,
                        "key_points": ["Main concepts", "Examples", "Applications"]
                    }
                fallback_structure.append(slide)
            
            return fallback_structure

class ContentBot:
    def generate_slide_content(self, slide_info, topic):
        prompt = f"""
    You are a professional presentation writer tasked with creating highly engaging and novel slide content. 
    Your goal is to transform the topic into compelling, audience-ready text that goes beyond generic explanations. 

    Topic: "{topic}"
    Slide Title: {slide_info['title']}
    Design Type: {slide_info['design_type']}
    Key Points to Cover: {', '.join(slide_info['key_points'])}

    Instructions:
    - Write 2–3 cohesive paragraphs that balance clarity with originality.
    - Present the information in a narrative style that feels fresh and thought-provoking, 
      avoiding clichés or textbook-like phrasing.
    - Use vivid analogies, real-world examples, or surprising insights to make the content memorable.
    - Maintain a professional, engaging tone that aligns with the design type 
      (e.g., formal for informational slides, visionary for conclusion slides, 
      energetic for introductory slides).
    - Ensure the content is concise enough for a presentation, yet rich in substance.
    - Focus on delivering value to the audience by highlighting practical relevance, 
      broader implications, or unexpected connections.
    - Return only the pure content text (no formatting, bullet points, or markdown).
    """
        
        try:
            response = llm.invoke(prompt)
            return response.content.strip()
        except Exception as e:
            print(f"Error generating slide content: {e}")
            return f"Content about {slide_info['title']} related to {topic}. This slide covers {', '.join(slide_info['key_points'])}."
    
    def generate_image_prompt(self, slide_info, topic):
        prompt = f"""
        Create a detailed image prompt for generating a professional image for a PowerPoint slide.
        
        Slide title: {slide_info['title']}
        Topic: {topic}
        Design type: {slide_info['design_type']}
        
        The image should be:
        - Professional and suitable for business presentations
        - High quality and visually appealing
        - Relevant to the slide topic
        - Clean and modern design
        - No text or words in the image
        
        Write a concise but descriptive prompt (max 50 words) for an AI image generator.
        Only return the prompt text, no quotes or extra formatting.
        """
        
        try:
            response = llm.invoke(prompt)
            return response.content.strip()
        except Exception as e:
            print(f"Error generating image prompt: {e}")
            return f"Professional illustration about {topic}, modern design, high quality, business presentation style"

# Initialize PPT Bot
ppt_bot = PPTBot()

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.get_json()
        topic = data.get('topic', '').strip()
        slides_count = int(data.get('slides_count', 5))
        
        if not topic:
            return jsonify({'success': False, 'error': 'Topic is required'})
        
        if slides_count < 3 or slides_count > 20:
            return jsonify({'success': False, 'error': 'Slides count must be between 3 and 20'})
        
        print(f"Starting presentation generation for topic: {topic}, slides: {slides_count}")
        
        # Generate presentation
        result = ppt_bot.generate_presentation(topic, slides_count)
        
        if result['success']:
            # Convert relative paths to URLs
            if result['slides_data']:
                for slide in result['slides_data']:
                    if slide['image_path']:
                        slide['image_url'] = url_for('static', filename=slide['image_path'])
            
            if result['ppt_path']:
                result['ppt_url'] = url_for('static', filename=result['ppt_path'])
        
        return jsonify(result)
        
    except Exception as e:
        print(f"Error in generate_ppt endpoint: {e}")
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': str(e),
            'message': 'Server error occurred'
        })

@app.route('/api/status')
def status():
    return jsonify({
        'status': 'running',
        'model_loaded': client is not None,
        'timestamp': datetime.now().isoformat(),
        'groq_available': GROQ_API_KEY is not None,
        'hf_available': HF_TOKEN is not None
    })

@app.route('/health')
def health():
    return jsonify({
        'status': 'healthy',
        'services': {
            'groq': GROQ_API_KEY is not None,
            'huggingface': client is not None,
            'directories_created': all([
                os.path.exists(UPLOAD_FOLDER),
                os.path.exists(IMAGES_FOLDER),
                os.path.exists(PPT_FOLDER)
            ])
        }
    })

if __name__ == '__main__':
    print("Starting AI PowerPoint Generator...")
    print(f"GROQ API Key: {'✓' if GROQ_API_KEY else '✗'}")
    print(f"HF Token: {'✓' if HF_TOKEN else '✗'}")
    print(f"Image Client: {'✓' if client else '✗'}")
    
    app.run(debug=True, host='0.0.0.0', port=5024)
